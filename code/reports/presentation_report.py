from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import sys
import os
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
from bs4 import BeautifulSoup
from utils.configs import llm
from PIL import Image
import plotly.graph_objects as go
from plotly.io import to_image
from markdown import markdown
from prompts.presentation_prompt_template import presentation_prompt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import random

def select_slide_layout(slide_content, include_plot):
    if 'section_title' in slide_content and 'content' in slide_content:
        if slide_content.get('table') is not None:
            return 4  # Content with Caption - TEXT & TABLE
        elif include_plot and slide_content.get('plot') is not None:
            layout_choice = random.choice([2, 3])  # Randomly choose between small and large plot layouts
            return layout_choice
        else:
            return 1  # Title and Content - TEXT
    elif 'report_title' in slide_content:
        return 0  # Title Slide - TITLE
    else:
        return 4  # Default layout

def get_presentation_content(section_content):
    content = section_content['content']
    prompt = presentation_prompt.replace("{section_content}", content)
    response = llm.get_response(prompt)
    return response

def parse_slides(section_content):
    response = get_presentation_content(section_content)
    soup = BeautifulSoup(markdown(response), 'html.parser')
    slides = []
    
    for slide in soup.find_all('slide'):
        title_element = slide.find('section_title')
        report_title_element = slide.find('report_title')
        report_title = report_title_element.text if report_title_element else None
        section_title = title_element.text if title_element else "Untitled Section"
        
        content = []
        table_text = ""
        in_table = False
        content_element = slide.find('content')
        if content_element:
            for line in content_element.text.strip().split('\n'):
                if line.strip().startswith('|'):
                    in_table = True
                    table_text += line + '\n'
                elif in_table:
                    in_table = False
                    table = parse_markdown_table(table_text)
                    table_text = ""
                elif line.strip():
                    if line.strip().startswith('- '):
                        content.append({'type': 'bullet', 'text': line.strip()[2:]})
                    else:
                        content.append({'type': 'paragraph', 'text': line.strip()})
        
        if table_text:
            table = parse_markdown_table(table_text)
        else:
            table = None
        
        plot = section_content.get('plot', None)
        
        slides.append({
            'report_title': report_title,
            'section_title': section_title,
            'content': content,
            'table': table,
            'plot': plot
        })
    
    return slides

def parse_markdown_table(markdown_table):
    lines = markdown_table.strip().split('\n')
    if len(lines) < 2:
        return None

    headers = [cell.strip() for cell in lines[0].split('|') if cell.strip()]
    data = []
    for line in lines[2:]:
        cells = [cell.strip() for cell in line.split('|') if cell.strip()]
        if cells:
            data.append(cells)
    
    if headers and data:
        return {'headers': headers, 'data': data}
    else:
        return None
    
def create_presentation(section_content, prs=None, selected_template='default'):
    if prs is None:
        if selected_template == 'default':
            prs = Presentation('code/templates/BlueYellow.pptx')
        prs = Presentation(f"code/templates/{selected_template}.pptx")
        
    # Add title slide only if it's the first section
    if not prs.slides:
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        report_title = title_slide.shapes.title
        report_title.text = section_content.get('report_title', 'Untitled Presentation')
    
    slides = parse_slides(section_content)
    total_slides = len(prs.slides) + len(slides)
    
    plot_added = False
    
    for slide_number, slide_content in enumerate(slides, start=len(prs.slides) + 1):
        include_plot = slide_content.get('plot') is not None and not plot_added
        layout_index = select_slide_layout(slide_content, include_plot)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        section_title = slide.shapes.title
        if section_title:
            section_title.text = slide_content.get('section_title', 'Untitled Section')
        
        if 'content' in slide_content:
            add_content_to_slide(slide, slide_content)
        
        if slide_content.get('table') is not None:
            try:
                add_table_to_slide(slide, slide_content['table'])
            except KeyError as e:
                print(f"Error adding table to slide: {e}")
        
        if include_plot:
            try:
                add_plot_to_slide(slide, slide_content['plot'])
                plot_added = True
            except KeyError as e:
                print(f"Error adding plot to slide: {e}")
    
    return prs

def add_content_to_slide(slide, slide_content):
    content_placeholder = find_content_placeholder(slide)
    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear any existing content

    for item in slide_content['content']:
        if item['type'] == 'paragraph':
            p = text_frame.add_paragraph()
            p.text = item['text']
            p.level = 0
        elif item['type'] == 'bullet':
            p = text_frame.add_paragraph()
            p.text = item['text']
            p.level = 1

def add_plot_to_slide(slide, plot_data):
    if not plot_data or 'data' not in plot_data or 'layout' not in plot_data:
        print("Invalid plot data")
        return

    try:
        plot_placeholder = find_plot_placeholder(slide)

        fig = go.Figure(data=plot_data['data'], layout=plot_data['layout'])
        img_bytes = to_image(fig, format="png", engine="kaleido", width=900, height=500, scale=2)

        plot_placeholder.insert_picture(BytesIO(img_bytes))

    except Exception as e:
        print(f"Error adding plot to slide: {e}")

    
def find_content_placeholder(slide):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx != 0:  # Skip title placeholder
            return placeholder
    raise KeyError("No content placeholder found on this slide")

def find_plot_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 18:  
            return shape
    raise KeyError("No plot placeholder found on this slide")

def find_table_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 12:  # Table placeholder
            return shape
    raise KeyError("No table placeholder found on this slide")

def add_summary_slide(prs, section_results):
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = summary_slide.shapes.title
    title.text = "Key Takeaways"

    content = summary_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for i, (name, _) in enumerate(section_results, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {name}"
        p.level = 0
        p.font.size = Pt(18)

    return prs

def add_table_to_slide(slide, table_data):
    if table_data is None or 'headers' not in table_data or 'data' not in table_data:
        print("Invalid table data")
        return
    try:
        table_placeholder = find_table_placeholder(slide)
        rows = len(table_data['data']) + 1  # +1 for header
        cols = len(table_data['headers'])
        
        table = table_placeholder.insert_table(rows, cols).table
        for col_idx, header in enumerate(table_data['headers']):
            table.cell(0, col_idx).text = header
        
        for row_idx, row in enumerate(table_data['data']):
            for col_idx, cell in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = cell
    except Exception as e:
        print(f"Error adding table to slide: {e}")